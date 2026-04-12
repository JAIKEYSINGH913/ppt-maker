#!/usr/bin/env python3
"""Print slide counts and shape counts per slide for reference PPTX (Sample Files)."""
from __future__ import annotations

import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("Install md2deck deps: pip install python-pptx")
    sys.exit(1)


def analyze_pptx(path: Path) -> None:
    if not path.exists():
        print(f"Missing: {path}")
        return
    prs = Presentation(str(path))
    print(f"\n=== {path.name} ===")
    print(f"slides: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides):
        try:
            layout = slide.slide_layout.name
        except Exception:
            layout = "?"
        n = len(slide.shapes)
        print(f"  slide {i+1:2d}  layout={layout!r:40s}  shapes={n}")


def main() -> None:
    root = Path(__file__).resolve().parents[2]
    samples = root / "Sample Files"
    paths = list(samples.rglob("*.pptx")) if samples.is_dir() else []
    paths = [p for p in paths if not p.name.startswith("~$")]
    if not paths:
        print("No PPTX under Sample Files")
        return
    for p in sorted(paths, key=lambda x: x.name.lower()):
        analyze_pptx(p)


if __name__ == "__main__":
    main()
