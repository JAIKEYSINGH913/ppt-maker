"""
Slide thumbnail generator for Spectral Weaver.
Generates PNG thumbnails from PPTX slide masters using PowerPoint COM on Windows.
Falls back to python-pptx + Pillow basic rendering if COM is unavailable.
"""
from __future__ import annotations

import sys
import subprocess
from pathlib import Path
from typing import List, Optional


def generate_thumbnails(
    pptx_path: Path,
    output_dir: Path,
    width: int = 960,
    height: int = 540,
    force: bool = False,
    as_layouts: bool = True,
) -> List[Path]:
    """Generate PNG thumbnails for each slide (or layout if as_layouts=True) in a PPTX file.

    Returns a sorted list of PNG paths.
    """
    output_dir.mkdir(parents=True, exist_ok=True)

    # Return cached thumbnails if they exist
    if not force:
        existing = sorted(output_dir.glob("slide_*.png"))
        if existing:
            return existing

    # Try PowerPoint COM first (best quality, Windows only)
    result = _try_powerpoint_com(pptx_path, output_dir, width, height, as_layouts)
    if result:
        return result

    # Try LibreOffice headless + pdf2image (Note: as_layouts not supported here yet)
    result = _try_libreoffice(pptx_path, output_dir)
    if result:
        return result

    # Last resort: basic Pillow rendering from python-pptx shapes
    result = _try_pillow_fallback(pptx_path, output_dir, width, height, as_layouts)
    if result:
        return result

    return []


# ────────────────────────  PowerPoint COM  ────────────────────────

def _try_powerpoint_com(
    pptx_path: Path, output_dir: Path, width: int, height: int, as_layouts: bool = False
) -> Optional[List[Path]]:
    if sys.platform != "win32":
        return None
    try:
        import comtypes.client  # type: ignore

        import comtypes
        comtypes.CoInitialize()
        ppt = comtypes.client.CreateObject("PowerPoint.Application")
        abs_path = str(pptx_path.resolve())
        prs = ppt.Presentations.Open(abs_path, ReadOnly=True, WithWindow=False)

        results: List[Path] = []
        
        if as_layouts:
            # Generate one slide per layout to get its background/placeholders
            # We use the existing presentation as it already contains the masters
            for idx in range(1, prs.SlideMaster.CustomLayouts.Count + 1):
                layout = prs.SlideMaster.CustomLayouts(idx)
                # Create a temporary slide with this layout
                temp_slide = prs.Slides.AddSlide(prs.Slides.Count + 1, layout)
                out = output_dir / f"slide_{idx - 1}.png"
                temp_slide.Export(str(out.resolve()), "PNG", width, height)
                # Delete the temporary slide immediately
                temp_slide.Delete()
                results.append(out)
        else:
            for idx in range(1, prs.Slides.Count + 1):
                out = output_dir / f"slide_{idx - 1}.png"
                prs.Slides(idx).Export(str(out.resolve()), "PNG", width, height)
                results.append(out)

        prs.Close()
        ppt.Quit()
        comtypes.CoUninitialize()
        print(f"[Thumbnailer] PowerPoint COM ({'Layouts' if as_layouts else 'Slides'}): {len(results)} generated")
        return results
    except Exception as exc:
        print(f"[Thumbnailer] PowerPoint COM failed: {str(exc)}")
        return None


# ────────────────────────  LibreOffice  ───────────────────────────

def _try_libreoffice(
    pptx_path: Path, output_dir: Path
) -> Optional[List[Path]]:
    soffice = _find_soffice()
    if not soffice:
        return None
    try:
        pdf_dir = output_dir / "_pdf"
        pdf_dir.mkdir(exist_ok=True)

        subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf",
             "--outdir", str(pdf_dir), str(pptx_path)],
            check=True, timeout=120, capture_output=True,
        )

        pdf_files = list(pdf_dir.glob("*.pdf"))
        if not pdf_files:
            return None

        from pdf2image import convert_from_path  # type: ignore
        images = convert_from_path(str(pdf_files[0]), dpi=150)

        results: List[Path] = []
        for i, img in enumerate(images):
            out = output_dir / f"slide_{i}.png"
            img.save(str(out), "PNG")
            results.append(out)

        print(f"[Thumbnailer] LibreOffice: {len(results)} slides -> {output_dir}")
        return results
    except Exception as exc:
        print(f"[Thumbnailer] LibreOffice failed: {exc}")
        return None


def _find_soffice() -> Optional[str]:
    candidates = [
        "soffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]
    for cmd in candidates:
        try:
            subprocess.run([cmd, "--version"], capture_output=True, timeout=5)
            return cmd
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue
    return None


# ────────────────────────  Pillow fallback  ───────────────────────

def _try_pillow_fallback(
    pptx_path: Path, output_dir: Path, width: int, height: int, as_layouts: bool = False
) -> Optional[List[Path]]:
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Emu          # type: ignore
        from PIL import Image, ImageDraw, ImageFont  # type: ignore

        prs = Presentation(str(pptx_path))
        sw = prs.slide_width or Emu(12192000)
        sh = prs.slide_height or Emu(6858000)
        scale_x = width / sw
        scale_y = height / sh

        results: List[Path] = []
        items = prs.slide_layouts if as_layouts else prs.slides
        
        for i, item in enumerate(items):
            img = Image.new("RGB", (width, height), (18, 18, 24))
            draw = ImageDraw.Draw(img)

            # For layouts, background is in the layout itself or the master
            try:
                bg = item.background
                if bg.fill and bg.fill.type is not None:
                    rgb = bg.fill.fore_color.rgb
                    img = Image.new("RGB", (width, height), tuple(rgb))
                    draw = ImageDraw.Draw(img)
            except Exception:
                pass

            # Render shapes
            for shape in item.shapes:
                try:
                    left = int((shape.left or 0) * scale_x)
                    top = int((shape.top or 0) * scale_y)
                    w = int((shape.width or 0) * scale_x)
                    h = int((shape.height or 0) * scale_y)

                    draw.rectangle(
                        [left, top, left + w, top + h],
                        outline=(80, 80, 100),
                    )

                    if shape.has_text_frame:
                        text = shape.text_frame.text[:100].strip()
                        if text:
                            try:
                                font = ImageFont.truetype("arial.ttf", max(10, min(h // 4, 18)))
                            except OSError:
                                font = ImageFont.load_default()
                            draw.text(
                                (left + 6, top + 6), text,
                                fill=(220, 220, 230), font=font,
                            )
                except Exception:
                    continue

            # Item label
            try:
                font_s = ImageFont.truetype("arial.ttf", 14)
            except OSError:
                font_s = ImageFont.load_default()
            label = f"L#{i}" if as_layouts else f"S#{i + 1}"
            draw.text((width - 50, height - 30), label, fill=(120, 120, 140), font=font_s)

            out = output_dir / f"slide_{i}.png"
            img.save(str(out), "PNG")
            results.append(out)

        print(f"[Thumbnailer] Pillow fallback ({'Layouts' if as_layouts else 'Slides'}): {len(results)} generated")
        return results
    except Exception as exc:
        print(f"[Thumbnailer] Pillow fallback failed: {str(exc)}")
        return None
