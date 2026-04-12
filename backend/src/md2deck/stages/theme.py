from __future__ import annotations

from dataclasses import dataclass

from pptx import Presentation

from md2deck.config import AppConfig
from md2deck.models import PipelineArtifacts, ThemeProfile, LayoutMetadata, SlotMetadata, SlideContentMetadata
from md2deck.stages.master_profiles import geometry_for_master_name


@dataclass(slots=True)
class ThemeStage:
    name: str = "theme"

    def run(self, config: AppConfig, artifacts: PipelineArtifacts) -> None:
        presentation = Presentation(str(config.master_pptx))
        master_key = config.master_pptx.name.lower()
        palette = self._palette_for_master(master_key)
        geometry = geometry_for_master_name(master_key)
        
        # Analyze all layouts
        layouts_metadata = {}
        from pptx.enum.shapes import PP_PLACEHOLDER
        
        for idx, layout in enumerate(presentation.slide_layouts):
            meta = LayoutMetadata(name=layout.name.lower(), index=idx)
            
            content_lefts, content_tops, content_rights, content_bottoms = [], [], [], []
            
            for ph in layout.placeholders:
                pt = ph.placeholder_format.type
                slot = SlotMetadata(
                    type=str(pt),
                    left=ph.left,
                    top=ph.top,
                    width=ph.width,
                    height=ph.height
                )
                
                if pt in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                    meta.has_title = True
                elif pt == PP_PLACEHOLDER.SUBTITLE:
                    meta.has_subtitle = True
                elif pt in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                    meta.body_slots.append(slot)
                    content_lefts.append(ph.left)
                    content_tops.append(ph.top)
                    content_rights.append(ph.left + ph.width)
                    content_bottoms.append(ph.top + ph.height)
                elif pt == PP_PLACEHOLDER.PICTURE:
                    meta.picture_slots.append(slot)
                    content_lefts.append(ph.left)
                    content_tops.append(ph.top)
                    content_rights.append(ph.left + ph.width)
                    content_bottoms.append(ph.top + ph.height)
            
            # Calculate Safe Zone bounding box of all content slots
            if content_lefts:
                meta.safe_rect = (
                    min(content_lefts),
                    min(content_tops),
                    max(content_rights) - min(content_lefts),
                    max(content_bottoms) - min(content_tops)
                )
            
            layouts_metadata[idx] = meta

        # Analyze existing slides for Semantic Injection
        template_dna = []
        for i, slide in enumerate(presentation.slides):
            texts = []
            tables = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t: texts.append(t)
                if shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        table_data.append([cell.text_frame.text.strip() for cell in row.cells])
                    tables.append(table_data)
            
            template_dna.append(SlideContentMetadata(
                index=i,
                layout_name=slide.slide_layout.name,
                texts=texts,
                tables=tables,
                shape_count=len(slide.shapes)
            ))

        artifacts.theme = ThemeProfile(
            master_path=config.master_pptx,
            slide_width=presentation.slide_width,
            slide_height=presentation.slide_height,
            layout_names=[layout.name for layout in presentation.slide_layouts],
            layouts_metadata=layouts_metadata,
            primary_color=palette["primary"],
            accent_colors=palette["accents"],
            light_colors=palette["lights"],
            dark_color=palette["dark"],
            muted_color=palette["muted"],
            theme_notes={
                "master_slide_count": len(presentation.slides),
                "visual_goal": "inherit master design while enabling infographic-first dynamic layouts",
                "analyzed_layouts": len(layouts_metadata),
            },
            design_dna=palette.get("dna") or {},
            geometry=geometry,
            template_dna=template_dna,
        )

    @staticmethod
    def _palette_for_master(master_name: str) -> dict[str, tuple | list | dict]:
        # Accenture - High Fidelity Detected Palette
        acc_purple = (127, 63, 255)    # #7F3FFF
        acc_cyan_purple = (162, 89, 255) # #A259FF
        acc_dark = (26, 26, 26)        # #1A1A1A
        acc_gray = (111, 111, 111)      # #6F6F6F
        acc_meta = (154, 154, 154)      # #9A9A9A
        white = (255, 255, 255)
        
        # AI Bubble - High Fidelity Detected Palette
        ai_navy = (11, 31, 58)         # #0B1F3A
        ai_blue = (47, 128, 237)        # #2F80ED
        ai_glow = (86, 204, 242)        # #56CCF2
        ai_gray = (176, 183, 195)       # #B0B7C3
        ai_text = (214, 226, 240)       # #D6E2F0
        
        # UAE Solar - High Fidelity Detected Palette
        uae_green = (11, 61, 46)        # #0B3D2E
        uae_teal = (26, 188, 156)       # #1ABC9C
        uae_yellow = (242, 201, 76)      # #F2C94C
        uae_gray = (95, 95, 95)         # #5F5F5F
        
        if "accenture" in master_name:
            return {
                "primary": acc_purple,
                "accents": [acc_purple, acc_cyan_purple, acc_gray],
                "lights": [white, (249, 249, 249)],
                "dark": acc_dark,
                "muted": acc_gray,
                "dna": {
                    "margin": 0.56,
                    "align": "center",
                    "title": {"h1": acc_purple, "h2": acc_gray, "meta": acc_meta},
                    "content": {"h1": acc_dark, "h2": acc_cyan_purple, "b": (51, 51, 51), "highlight": acc_purple},
                    "closing": {"h1": acc_purple, "h2": acc_gray}
                }
            }
            
        if "ai bubble" in master_name or "bubble" in master_name:
            return {
                "primary": ai_navy,
                "accents": [ai_glow, ai_blue, ai_gray],
                "lights": [ai_text, white],
                "dark": ai_navy,
                "muted": ai_gray,
                "dna": {
                    "margin": 0.38,
                    "align": "left",
                    "title": {"h1": ai_glow, "h2": ai_gray},
                    "content": {"h1": white, "h2": ai_blue, "b": ai_text, "highlight": ai_glow},
                    "closing": {"h1": ai_glow, "h2": ai_gray}
                }
            }
            
        if "uae" in master_name or "solar" in master_name:
            uae_green_dna = (33, 62, 16) # #213E10 - extracted Green
            return {
                "primary": uae_green_dna,
                "accents": [uae_yellow, uae_teal, uae_green_dna],
                "lights": [white, (230, 245, 240)],
                "dark": uae_green_dna,
                "muted": uae_gray,
                "dna": {
                    "margin": 0.43,
                    "align": "left",
                    "title": {"h1": uae_green_dna, "h2": uae_gray},
                    "content": {"h1": uae_green_dna, "h2": uae_teal, "b": (34, 34, 34), "highlight": uae_yellow},
                    "closing": {"h1": uae_green_dna, "h2": uae_yellow}
                }
            }

        return {
            "primary": (15, 23, 42),
            "accents": [(15, 23, 42), (14, 165, 233), (16, 185, 129)],
            "lights": [(241, 245, 249)],
            "dark": (15, 23, 42),
            "muted": (100, 116, 139),
            "dna": {}
        }
