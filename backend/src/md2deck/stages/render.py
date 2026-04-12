"""
Slide Renderer — Spectral Studio v3 (Safe-Stack Architecture)
=============================================================
Layout rules (from sample analysis):
  Slide: 13.333" x 7.5"
  Accenture margins: L=0.375", T=0.656" (title zone), content starts ~1.25"
  AI Bubble margins:  L=0.375", title at T=0.656"
  UAE Solar:          L=0.426", title at T=0.42", subtitle at T=1.125"

Font hierarchy (user-specified):
  Hero heading  : 40pt bold
  Slide title   : 25pt bold
  Subtitle      : 16pt
  Body/bullets  : 10pt
  Table header  : 16pt bold
  Table body    : 11pt
  Footer/number : 9pt

Safe-Stack principle:
  - Each _add_stacked_box() call advances self._stack_y by the box height + gap
  - Nothing is ever placed above self._stack_y
  - Hard clip: no box bottom may exceed SLIDE_H - 0.15"
  - Internal padding: 0.1" top/bottom, 0.12" left/right on every text frame
"""
from __future__ import annotations

import logging
import re
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt, Emu

from md2deck.config import AppConfig
from md2deck.models import (
    MasterLayoutProfile,
    PipelineArtifacts,
    SlideBlueprint,
    VisualIntent,
)
from md2deck.stages.slide_placeholders import fill_divider_placeholders, fill_title_subtitle_placeholders
from md2deck.stages.layout_engine import (
    GridLayout, Rect, clamp_text, clamp_lines,
    MARGIN_LEFT, MARGIN_RIGHT, MARGIN_TOP,
    CONTENT_TOP, CONTENT_BOTTOM, CONTENT_LEFT, CONTENT_WIDTH, CONTENT_HEIGHT,
    SLIDE_W, SLIDE_H,
    SHAPE_PAD_LEFT, SHAPE_PAD_RIGHT, SHAPE_PAD_TOP, SHAPE_PAD_BOTTOM,
)
from md2deck.stages.freepik_icons import fetch_icons_for_tokens, get_unicode_icon
from md2deck.stages.pptx_animations import add_slide_animations

logger = logging.getLogger(__name__)

# ── Slide constants ────────────────────────────────────────────────────────────
SW = Inches(13.333)   # Slide width
SH = Inches(7.5)      # Slide height
BOTTOM_SAFE = SH - Inches(0.2)   # Strictest bottom boundary
FOOTER_H = Inches(0.55)           # Storyteller footer height
FOOTER_T = SH - Inches(0.08) - FOOTER_H  # Footer top anchor

# ── Font sizes ─────────────────────────────────────────────────────────────────
FS_HERO   = 36   # Hero/cover headings (body slides use geometry.hero_title_pt on cover)
FS_TITLE  = 22   # Slide titles on Blank layouts
FS_SUB    = 16   # Subtitles, table headers
FS_BODY   = 10   # Default body/bullets
FS_FOOTER =  9   # Slide numbers, footer notes

# ── Internal text box padding ──────────────────────────────────────────────────
PAD_H = Inches(0.1)    # Top & bottom internal padding
PAD_W = Inches(0.12)   # Left & right internal padding

# ── Gap between stacked boxes ──────────────────────────────────────────────────
GAP_AFTER_TITLE = Inches(0.08)
GAP_AFTER_SUB   = Inches(0.12)
GAP_AFTER_BODY  = Inches(0.10)

MAX_CHARS_TITLE = 60


@dataclass(slots=True)
class RenderStage:
    name: str = "render"
    theme: object | None = None
    _stack_y: int = 0          # Current vertical stack cursor (EMU)
    _slide_margin: float = 0.5  # inches fallback when geometry frame not used
    _draw_frame: tuple[int, int, int, int] | None = None  # (left, top, width, height) EMU for Blank slides
    _icon_right_pad_emu: int = 0  # shrink text column when corner icon present

    # ──────────────────────────────────────────────────────────────────────────
    def run(self, config: AppConfig, artifacts: PipelineArtifacts) -> None:
        if artifacts.blueprint is None or artifacts.theme is None:
            raise RuntimeError("Blueprint and theme stages must run before rendering.")

        self.theme = artifacts.theme
        presentation = Presentation(str(config.master_pptx))
        presentation.slide_width = SW
        presentation.slide_height = SH

        # Keep a reference to the original presentation for semantic cloning
        source_prs = Presentation(str(config.master_pptx))

        self._trim_sample_slides(presentation)

        icon_cache_dir = config.working_dir / "icon_cache"

        total = len(artifacts.blueprint.slides)
        # Strategy: We will use the master PPTX itself as the base to preserve all formatting/logos.
        # We then duplicate the source slides mentioned in the blueprint.
        
        # 1. First, identify which template slides are needed
        final_slides_data = []
        for slide_idx, blueprint in enumerate(artifacts.blueprint.slides):
            role = self._slide_role(slide_idx, total, blueprint)
            
            # Map to template slide if index exists
            source_slide = None
            if blueprint.source_slide_index is not None and blueprint.source_slide_index < len(source_prs.slides):
                source_slide = source_prs.slides[blueprint.source_slide_index]
            
            final_slides_data.append((blueprint, source_slide, role))

        # 2. Clear all existing slides in the main presentation to start fresh but with the correct Master
        self._trim_sample_slides(presentation)

        # 3. Create the deck by either cloning source slides or adding layouts
        for blueprint, source_slide, role in final_slides_data:
            if source_slide:
                # To clone perfectly, we find the layout of the source and add it, then clone shapes.
                layout = source_slide.slide_layout
                slide = presentation.slides.add_slide(layout)
                self._clone_slide_content(source_slide, slide)
                self._inject_semantic_content(slide, blueprint)
            else:
                layout = self._resolve_layout(presentation, role)
                slide = presentation.slides.add_slide(layout)
                self._render_slide_for_role(slide, blueprint, 0, total, role)

            # Freepik icons (skip on hero to avoid clashing with master brand art)
            icon_paths: list[str] = []
            if role != "hero" and config.freepik_api_key and blueprint.icon_tokens:
                try:
                    icon_paths = [
                        str(p) for p in fetch_icons_for_tokens(
                            config.freepik_api_key,
                            blueprint.icon_tokens[:3],
                            icon_cache_dir,
                        )
                    ]
                except Exception as e:
                    logger.warning(f"Icon fetch failed slide {slide_idx}: {e}")
            blueprint.icon_paths = icon_paths

            self._draw_frame = None
            self._icon_right_pad_emu = 0
            if role == "content":
                self._draw_frame = self._compute_blank_content_frame(presentation)

            if blueprint.source_slide_index is not None:
                self._inject_semantic_content(slide, blueprint)
                # Ensure contrast/layering is okay after injection
            else:
                self._render_slide_for_role(slide, blueprint, slide_idx, total, role)

            if config.animations.transitions_enabled and role != "hero":
                try:
                    add_slide_animations(
                        slide,
                        blueprint.visual_intent.value,
                        len(slide.shapes),
                        include_entrance=config.animations.entrance_animations_enabled,
                    )
                except Exception as e:
                    logger.debug("Slide animation skipped: %s", e)

        presentation.save(str(config.output_pptx))

    # ──────────────────────────────────────────────────────────────────────────
    @staticmethod
    def _trim_sample_slides(presentation: Presentation) -> None:
        while len(presentation.slides) > 0:
            sid = presentation.slides._sldIdLst[0]
            presentation.part.drop_rel(sid.rId)
            del presentation.slides._sldIdLst[0]

    # ──────────────────────────────────────────────────────────────────────────
    @staticmethod
    def _slide_role(idx: int, total: int, blueprint: SlideBlueprint) -> str:
        if idx == 0:
            return "hero"
        if idx == total - 1:
            return "thank_you"
        # Only true divider slides (explicit hint); SECTION_DIVIDER intent still uses Blank + body.
        if str(blueprint.layout_hint or "").lower() == "divider":
            return "divider"
        return "content"

    @staticmethod
    def _resolve_layout(prs: Presentation, role: str):
        names = {ly.name.lower(): ly for ly in prs.slide_layouts}

        if role == "hero":
            for key in ("0_title", "1_cover", "2_cover", "cover", "title company", "title"):
                for n, ly in names.items():
                    if key in n and "blank" not in n:
                        return ly
            return prs.slide_layouts[0]

        if role == "thank_you":
            for key in ("thank you", "thankyou", "1_thank", "closing", "end", "done"):
                for n, ly in names.items():
                    if key in n:
                        return ly
            return prs.slide_layouts[-1]

        if role == "divider":
            for key in ("divider", "section header", "section", "c_section blue", "c_section"):
                if key in names:
                    return names[key]
            if "title only" in names:
                return names["title only"]
            return prs.slide_layouts[0]

        if "blank" in names:
            return names["blank"]
        return prs.slide_layouts[min(3, len(prs.slide_layouts) - 1)]

    def _default_geometry(self) -> MasterLayoutProfile:
        return self.theme.geometry or MasterLayoutProfile()

    def _compute_blank_content_frame(self, prs: Presentation) -> tuple[int, int, int, int]:
        """Tighter of per-master insets and Blank layout safe_rect (EMU)."""
        names = [ly.name.lower() for ly in prs.slide_layouts]
        blank_idx = next((i for i, n in enumerate(names) if n == "blank"), min(3, len(names) - 1))
        meta = self.theme.layouts_metadata.get(blank_idx)
        geo = self._default_geometry()

        inset_l = int(Inches(geo.blank_inset_left))
        inset_t = int(Inches(geo.blank_inset_top))
        inset_r = int(Inches(geo.blank_inset_right))
        inset_b = int(Inches(geo.blank_inset_bottom)) + int(Inches(geo.footer_reserve_inches))

        l, t = inset_l, inset_t
        r, b = int(SW) - inset_r, int(SH) - inset_b

        if meta and meta.safe_rect:
            sl, st, sw, sh = meta.safe_rect
            sr, sb = sl + sw, st + sh
            l = max(l, sl)
            t = max(t, st)
            r = min(r, sr)
            b = min(b, sb)

        w = max(int(Inches(0.8)), r - l)
        h = max(int(Inches(0.8)), b - t)
        return l, t, w, h

    def _render_slide_for_role(
        self, slide, blueprint: SlideBlueprint, idx: int, total: int, role: str
    ) -> None:
        dna = self.theme.design_dna
        self._slide_margin = float(dna.get("margin", 0.5))
        brand_align = dna.get("align", "left")
        pp_align = PP_ALIGN.CENTER if brand_align == "center" else PP_ALIGN.LEFT
        geo = self._default_geometry()
        c_title = self._dna_color(idx, total, "h1")
        c_subtitle = self._dna_color(idx, total, "h2")

        if role == "hero":
            sub = blueprint.summary or (blueprint.data_points[0] if blueprint.data_points else "")
            fill_title_subtitle_placeholders(
                slide,
                title=blueprint.title or "",
                subtitle=sub or None,
                title_rgb=c_title,
                subtitle_rgb=c_subtitle,
                title_pt=geo.hero_title_pt,
                subtitle_pt=FS_SUB,
                title_align=pp_align,
                max_title_chars=geo.hero_title_max_chars,
            )
            return

        if role == "thank_you":
            fill_title_subtitle_placeholders(
                slide,
                title=blueprint.title or "Thank You",
                subtitle=blueprint.summary or None,
                title_rgb=c_title,
                subtitle_rgb=c_subtitle,
                title_pt=FS_TITLE + 2,
                subtitle_pt=FS_SUB,
                title_align=pp_align,
                max_title_chars=48,
            )
            self._add_slide_number(slide, idx + 1, total)
            return

        if role == "divider":
            fill_divider_placeholders(
                slide,
                title=blueprint.title or "",
                subtitle=blueprint.summary or None,
                title_rgb=c_title,
                subtitle_rgb=c_subtitle,
                title_pt=FS_TITLE,
                subtitle_pt=FS_SUB - 1,
                title_align=pp_align,
                max_title_chars=72,
            )
            self._add_slide_number(slide, idx + 1, total)
            return

    def _clone_slide_content(self, source_slide, dest_slide) -> None:
        """Wipe destination shapes and deep copy ALL shapes from source, including placeholders."""
        import copy
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        
        # 1. Wipe automatically created shapes from add_slide(layout) to make room for full clone
        for shape in list(dest_slide.shapes):
            try:
                sp = shape.element
                sp.getparent().remove(sp)
            except Exception:
                pass

        # 2. Clone every shape from source slide (logos, backdrops, text boxes)
        for shape in source_slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Proper image cloning with blob to preserve relationships
                from io import BytesIO
                img_stream = BytesIO(shape.image.blob)
                new_pic = dest_slide.shapes.add_picture(img_stream, shape.left, shape.top, shape.width, shape.height)
                # Ensure it keeps same name for potentially later lookup
                new_pic.name = shape.name
            else:
                # Copy by XML for vector shapes, lines, and text boxes
                new_sp = copy.deepcopy(shape.element)
                dest_slide.shapes._spTree.append(new_sp)

    def _copy_font_style(self, source_para, target_para) -> None:
        """Deep copy font properties (color, size, bold) from one paragraph to another."""
        try:
            sf = source_para.font
            tf = target_para.font
            if sf.name: tf.name = sf.name
            if sf.size: tf.size = sf.size
            if sf.bold is not None: tf.bold = sf.bold
            if sf.italic is not None: tf.italic = sf.italic
            if sf.color and hasattr(sf.color, 'rgb') and sf.color.rgb:
                tf.color.rgb = sf.color.rgb
        except Exception:
            pass

    def _rewrite_text_frame(self, tf, lines: list[str]) -> None:
        """Populate a text frame while attempting to preserve template font styling."""
        if not lines:
            tf.clear()
            return
            
        sample_para = None
        if tf.paragraphs:
            sample_para = tf.paragraphs[0]
            
        tf.clear()
        for line in lines:
            p = tf.add_paragraph()
            p.text = str(line)
            p.level = 0
            if sample_para:
                self._copy_font_style(sample_para, p)

    def _inject_semantic_content(self, slide, blueprint: SlideBlueprint) -> None:
        """Find existing text/tables/charts on a cloned slide and swap them with MD data."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
        from pptx.chart.data import CategoryChartData
        import urllib.parse

        # 1. Title handling
        title_box = None
        for s in slide.shapes:
            if s.is_placeholder and s.placeholder_format.type in (1, 3):
                title_box = s
                break
            if "title" in s.name.lower() and s.has_text_frame:
                title_box = s
                break
        if title_box and title_box.has_text_frame:
            self._rewrite_text_frame(title_box.text_frame, [blueprint.title])
            
        # 2. Table handling: Data mapping + Pruning
        tables = [s for s in slide.shapes if s.has_table]
        if tables and blueprint.table_spec:
            table = tables[0].table
            spec = blueprint.table_spec
            rows_data = spec.get("rows", [])
            header_data = spec.get("headers", [])
            
            # Map headers
            if header_data:
                for c_idx, h_val in enumerate(header_data[:len(table.columns)]):
                    try: table.cell(0, c_idx).text = str(h_val)
                    except: pass
            
            # Map rows
            for r_idx, row_data in enumerate(rows_data[:len(table.rows)-1]):
                for c_idx, cell_val in enumerate(row_data[:len(table.columns)]):
                    try: table.cell(r_idx + 1, c_idx).text = str(cell_val)
                    except: pass
            
            # Pruning logic: Clear extra columns
            if header_data and len(header_data) < len(table.columns):
                for c_idx in range(len(header_data), len(table.columns)):
                    for r_idx in range(len(table.rows)):
                        try: table.cell(r_idx, c_idx).text = ""
                        except: pass

        # 3. Chart handling
        charts = [s for s in slide.shapes if s.has_chart]
        if charts and blueprint.chart_spec:
            try:
                chart = charts[0].chart
                spec = blueprint.chart_spec
                chart_data = CategoryChartData()
                chart_data.categories = spec.get("categories", ["A", "B", "C"])
                series_list = spec.get("series", [])
                for s_info in series_list:
                    chart_data.add_series(s_info.get("name", "Value"), s_info.get("values", [1, 2, 3]))
                chart.replace_data(chart_data)
            except Exception as e:
                logger.warning(f"Chart injection failed: {e}")

        # 4. Image handling: Contextual replacement at same Z-Order
        pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE or (s.is_placeholder and s.placeholder_format.type == 14)]
        if pics and blueprint.title:
            pics.sort(key=lambda s: s.width * s.height, reverse=True)
            target_pic = pics[0]
            try:
                prompt = f"minimalist business concept illustration: {blueprint.title}"
                encoded = urllib.parse.quote(prompt)
                img_url = f"https://image.pollinations.ai/prompt/{encoded}?width=1024&height=768&nologo=true"
                from io import BytesIO
                import requests
                r = requests.get(img_url, timeout=10)
                if r.status_code == 200:
                    img_stream = BytesIO(r.content)
                    # Insert new picture at exact index of old picture to preserve layering
                    old_idx = list(slide.shapes._spTree).index(target_pic.element)
                    new_pic = slide.shapes.add_picture(img_stream, target_pic.left, target_pic.top, target_pic.width, target_pic.height)
                    slide.shapes._spTree.remove(new_pic.element)
                    slide.shapes._spTree.insert(old_idx, new_pic.element)
                    
                    sp = target_pic.element
                    sp.getparent().remove(sp)
            except Exception as e:
                logger.warning(f"Image replacement failed: {e}")

        # 5. Text Body & Summary Distribution (Preserving Font Style)
        text_shapes = [s for s in slide.shapes if s.has_text_frame and s != title_box]
        text_shapes.sort(key=lambda s: s.width * s.height, reverse=True)
        if text_shapes:
            if len(text_shapes) >= 2 and blueprint.summary:
                body_box, summary_box = text_shapes[0], text_shapes[1]
                self._rewrite_text_frame(summary_box.text_frame, [blueprint.summary])
                self._rewrite_text_frame(body_box.text_frame, [str(pt) for pt in blueprint.data_points])
            else:
                body_box = text_shapes[0]
                lines = [str(pt) for pt in blueprint.data_points] if blueprint.data_points else ([blueprint.summary] if blueprint.summary else [])
                self._rewrite_text_frame(body_box.text_frame, lines)

    def _render_content_slide(self, slide, blueprint: SlideBlueprint, idx: int, total: int) -> None:
        """Stacked title, body, charts on Blank layout using _draw_frame (no duplicate placeholders)."""
        if self._draw_frame is None:
            margin = Inches(self._slide_margin)
            self._draw_frame = (int(margin), int(margin), int(SW - 2 * margin), int(SH - 2 * margin - Inches(0.65)))

        # Find native placeholders to avoid overlapping template elements (logos, text boxes).
        title_ph_done = False
        native_body_frame = None
        to_delete = []

        for ph in slide.placeholders:
            try:
                pf = ph.placeholder_format
            except Exception:
                continue
            if pf.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                # Use this bounding box instead of our default hardcoded frame!
                # This guarantees we don't overlap with master template designs.
                native_body_frame = (ph.left, ph.top, ph.width, ph.height)
                to_delete.append(ph)
            elif pf.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                to_delete.append(ph) # To prevent overlap, we will draw our own title or we can leave it

        # Delete unused native placeholders so "Click to add text" ghosts don't appear / overlap
        for ph in to_delete:
            try:
                sp = ph.element
                sp.getparent().remove(sp)
            except Exception:
                pass

        if native_body_frame:
            self._draw_frame = native_body_frame

        fl, ft, fw, fh = self._draw_frame
        dna = self.theme.design_dna
        brand_align = dna.get("align", "left")
        pp_align = PP_ALIGN.CENTER if brand_align == "center" else PP_ALIGN.LEFT
        geo = self._default_geometry()
        meta = blueprint.user_overrides or {}

        pad = int(Inches(0.05))
        margin_l = fl + pad
        content_w = fw - 2 * pad
        icon_hint = str(meta.get("icon_hint", "") or "")
        if icon_hint:
            self._icon_right_pad_emu = int(Inches(0.52) + Inches(geo.icon_margin_right))
        else:
            self._icon_right_pad_emu = 0

        self._draw_accent_strip(slide, fl, ft, fw)
        bar_h = int(Inches(0.055))
        self._stack_y = ft + bar_h + int(Inches(0.04))
        is_closing = idx == total - 1

        c_title = self._dna_color(idx, total, "h1")
        c_subtitle = self._dna_color(idx, total, "h2")
        c_body = self._dna_color(idx, total, "b")

        title_text = clamp_text(blueprint.title or "", max_words=8, max_chars=MAX_CHARS_TITLE)
        self._stack_box(
            slide, title_text, FS_TITLE, c_title,
            bold=True, align=pp_align, l=margin_l, w=content_w - self._icon_right_pad_emu,
            gap_after=GAP_AFTER_TITLE,
        )

        if blueprint.summary and not is_closing:
            summary = clamp_text(blueprint.summary, max_words=16, max_chars=96)
            self._stack_box(
                slide, summary, FS_SUB, c_subtitle,
                align=pp_align, l=margin_l, w=content_w - self._icon_right_pad_emu,
                gap_after=GAP_AFTER_SUB,
            )

        body_top = self._stack_y
        frame_bottom = ft + fh
        body_h_max = frame_bottom - body_top - int(Inches(0.18))
        if body_h_max < Inches(0.45):
            logger.warning("Slide %s: no body space in frame, skip body", idx)
        else:
            two_col = (
                blueprint.visual_intent == VisualIntent.CHART_FOCUS
                and "two_col" in str(blueprint.layout_hint or "").lower()
            )

            if blueprint.visual_intent == VisualIntent.METRIC_DASHBOARD:
                card_h = min(body_h_max, Inches(2.45))
                self._add_metric_cards(
                    slide, blueprint.chart_spec, margin_l, body_top, content_w - self._icon_right_pad_emu, card_h
                )
                self._stack_y = body_top + card_h + Inches(0.12)

            elif two_col:
                split_h = min(body_h_max, Inches(3.1))
                self._add_two_col(
                    slide,
                    blueprint.data_points,
                    self._ensure_chart_spec(blueprint),
                    margin_l,
                    body_top,
                    content_w - self._icon_right_pad_emu,
                    split_h,
                )
                self._stack_y = body_top + split_h + Inches(0.1)

            elif blueprint.visual_intent == VisualIntent.CHART_FOCUS:
                chart_h = min(body_h_max, Inches(3.65))
                self._add_chart(
                    slide,
                    self._ensure_chart_spec(blueprint),
                    margin_l,
                    body_top,
                    content_w - self._icon_right_pad_emu,
                    chart_h,
                )
                self._stack_y = body_top + chart_h + Inches(0.1)

            elif blueprint.visual_intent in (
                VisualIntent.ICON_GRID,
                VisualIntent.INFOGRAPHIC,
            ):
                grid_h = min(body_h_max, Inches(2.65))
                self._add_icon_grid(
                    slide,
                    blueprint.data_points,
                    blueprint.icon_paths,
                    blueprint.icon_tokens,
                    margin_l,
                    body_top,
                    content_w - self._icon_right_pad_emu,
                    grid_h,
                )
                self._stack_y = body_top + grid_h + Inches(0.12)
                
            elif blueprint.visual_intent == VisualIntent.CHEVRON_FLOW:
                flow_h = min(body_h_max, Inches(2.5))
                stages = blueprint.data_points[:5] or ["Phase 1", "Phase 2", "Phase 3"]
                self._add_chevron_flow(
                    slide, stages, left=margin_l, top=body_top,
                    width=content_w - self._icon_right_pad_emu, height=flow_h,
                )
                self._stack_y = body_top + flow_h + Inches(0.12)

            elif blueprint.visual_intent in (VisualIntent.FUNNEL, VisualIntent.PYRAMID):
                fh2 = min(body_h_max, Inches(2.9))
                stages = blueprint.data_points[:4] or [blueprint.title, "Stage 2", "Stage 3"]
                if blueprint.visual_intent == VisualIntent.FUNNEL:
                    self._add_funnel_diagram(
                        slide, stages, left=margin_l, top=body_top,
                        width=content_w - self._icon_right_pad_emu, height=fh2,
                    )
                else:
                    self._add_pyramid_diagram(
                        slide, stages, left=margin_l, top=body_top,
                        width=content_w - self._icon_right_pad_emu, height=fh2,
                    )
                self._stack_y = body_top + fh2 + Inches(0.1)

            elif blueprint.visual_intent == VisualIntent.TIMELINE:
                line_h = min(body_h_max, Inches(1.75))
                self._add_timeline(
                    slide,
                    blueprint.data_points,
                    margin_l,
                    body_top,
                    content_w - self._icon_right_pad_emu,
                    line_h,
                )
                self._stack_y = body_top + line_h + Inches(0.12)

            elif blueprint.visual_intent == VisualIntent.TABLE_FOCUS and blueprint.table_spec:
                tbl_h = min(body_h_max, Inches(3.05))
                self._add_table(
                    slide,
                    blueprint.table_spec,
                    margin_l,
                    body_top,
                    content_w - self._icon_right_pad_emu,
                    tbl_h,
                )
                self._stack_y = body_top + tbl_h + Inches(0.1)

            elif blueprint.visual_intent == VisualIntent.SWOT and blueprint.data_points:
                quad_h = min(body_h_max, Inches(3.05))
                self._add_swot(
                    slide,
                    blueprint.data_points,
                    margin_l,
                    body_top,
                    content_w - self._icon_right_pad_emu,
                    quad_h,
                )
                self._stack_y = body_top + quad_h + Inches(0.1)

            elif blueprint.visual_intent == VisualIntent.AGENDA:
                agenda_h = min(body_h_max, Inches(3.5))
                self._add_agenda(
                    slide,
                    blueprint.data_points,
                    margin_l,
                    body_top,
                    content_w - self._icon_right_pad_emu,
                    agenda_h,
                )
                self._stack_y = body_top + agenda_h + Inches(0.1)

            elif blueprint.visual_intent == VisualIntent.KEY_TAKEAWAYS:
                take_h = min(body_h_max, Inches(3.2))
                self._add_key_takeaways(
                    slide,
                    blueprint.data_points,
                    margin_l,
                    body_top,
                    content_w - self._icon_right_pad_emu,
                    take_h,
                )
                self._stack_y = body_top + take_h + Inches(0.1)

            elif blueprint.data_points:
                is_table_layout = blueprint.visual_intent == VisualIntent.TABLE_FOCUS
                body_size = FS_SUB if is_table_layout else FS_BODY
                bold_body = is_table_layout
                lines = []
                for pt in blueprint.data_points[: min(5, 8)]:
                    line = clamp_text(str(pt), max_words=14, max_chars=76)
                    lines.append(f"\u2022  {line}")
                body_text = "\n".join(lines)
                
                # If slide is virtually empty, use colorful palette text
                if len(blueprint.data_points) <= 1 and len(body_text) < 50:
                    body_size = FS_SUB + 4
                    c_body = self._primary() # Colorful palette

                self._stack_box(
                    slide,
                    body_text,
                    body_size,
                    c_body,
                    bold=bold_body,
                    align=PP_ALIGN.LEFT,
                    l=margin_l + int(Inches(0.04)),
                    w=content_w - self._icon_right_pad_emu - int(Inches(0.08)),
                    gap_after=GAP_AFTER_BODY,
                    hard_max_h=body_h_max,
                )

        if blueprint.canva_thumb_path and Path(blueprint.canva_thumb_path).exists():
            try:
                # If there are NO major points, embed large image in center instead of a tiny thumbnail
                if not blueprint.data_points or (len(blueprint.data_points) == 1 and len(blueprint.data_points[0]) < 20):
                    tw, th = fw - int(Inches(0.4)), fh - int(Inches(1.0))
                    slide.shapes.add_picture(
                        blueprint.canva_thumb_path,
                        fl + int(Inches(0.2)),
                        self._stack_y + int(Inches(0.1)),
                        width=tw,
                        height=th,
                    )
                else:    
                    tw, th = int(Inches(1.38)), int(Inches(0.78))
                    slide.shapes.add_picture(
                        blueprint.canva_thumb_path,
                        fl + fw - tw - int(Inches(0.06)),
                        ft + fh - th - int(Inches(0.12)),
                        width=tw,
                        height=th,
                    )
            except Exception as e:
                logger.debug("Canva thumb skip: %s", e)

        if icon_hint:
            icon_x = fl + fw - int(Inches(0.5)) - int(Inches(geo.icon_margin_right))
            icon_y = ft + int(Inches(geo.icon_margin_top))
            self._add_icon(slide, icon_hint, icon_x, icon_y)
        elif blueprint.icon_tokens and not self._any_local_icon_paths(blueprint.icon_paths):
            self._draw_unicode_token_row(
                slide, blueprint.icon_tokens[:3], fl + fw - int(Inches(2.0)), ft + int(Inches(0.1))
            )

        insight = meta.get("storyteller_insight", "")
        if insight and not is_closing:
            self._add_footer_box(slide, insight, margin_l, content_w - self._icon_right_pad_emu)

        self._add_slide_number(slide, idx + 1, total)

    @staticmethod
    def _any_local_icon_paths(paths: list[str]) -> bool:
        return any(p and Path(p).exists() for p in (paths or []))

    def _draw_accent_strip(self, slide, fl: int, ft: int, fw: int) -> None:
        """Thin brand accent bar at top of content frame."""
        bar_h = int(Inches(0.055))
        try:
            bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, fl, ft, fw, bar_h)
            bar.fill.solid()
            bar.fill.fore_color.rgb = self._accent(0)
            try:
                bar.line.fill.background()
            except Exception:
                pass
        except Exception as e:
            logger.debug("Accent strip skip: %s", e)

    def _draw_unicode_token_row(self, slide, tokens: list[str], left: int, top: int) -> None:
        """Small Unicode glyphs when raster icons are unavailable."""
        x = left
        for tok in tokens:
            ch = get_unicode_icon(tok)
            box = slide.shapes.add_textbox(x, top, Inches(0.35), Inches(0.35))
            p = box.text_frame.paragraphs[0]
            p.text = ch
            p.font.size = Pt(20)
            p.font.color.rgb = self._primary()
            p.alignment = PP_ALIGN.CENTER
            x += int(Inches(0.38))

    def _ensure_chart_spec(self, blueprint: SlideBlueprint) -> dict | None:
        spec = blueprint.chart_spec
        if spec and spec.get("labels"):
            vals = spec.get("values") or []
            if vals and any(float(v) != 0.0 for v in vals):
                return spec
        pairs = self._extract_chart_pairs(list(blueprint.data_points or []))
        if len(pairs) >= 2:
            return {"labels": [p[0] for p in pairs], "values": [p[1] for p in pairs]}
        return spec

    # ──────────────────────────────────────────────────────────────────────────
    # SAFE STACK BOX
    # ──────────────────────────────────────────────────────────────────────────
    def _stack_box(self, slide, text: str, size: int, color: RGBColor,
                   *, bold: bool = False, align=PP_ALIGN.LEFT,
                   l: int = None, w: int = None,
                   gap_after: int = 0, hard_max_h: int = 0) -> None:
        """Add a text box below the current stack cursor, advance cursor."""
        if not text or not text.strip():
            return

        margin = Inches(self._slide_margin)
        if self._draw_frame is not None:
            fl, ft, fw, fh = self._draw_frame
            default_l = fl + int(Inches(0.04))
            default_w = fw - int(Inches(0.08)) - self._icon_right_pad_emu
            bottom_cap = ft + fh - int(Inches(0.06))
        else:
            default_l = margin
            default_w = SW - 2 * margin
            bottom_cap = int(BOTTOM_SAFE)

        l = l if l is not None else default_l
        w = w if w is not None else default_w

        # Estimate height (slightly conservative line height to reduce overflow)
        usable_w_pts = ((w - 2 * PAD_W) / 914400) * 72  # inches -> points
        chars_per_line = max(1, int(usable_w_pts / (size * 0.58)))
        num_lines = 0
        for line in text.split("\n"):
            num_lines += max(1, -(-len(line) // chars_per_line))  # ceiling division

        line_h_emu = int(size * 1.58 * 12700)
        est_h = num_lines * line_h_emu + int(2 * PAD_H)

        # Apply hard max constraint
        if hard_max_h > 0 and est_h > hard_max_h:
            # Shrink: recalculate how many lines fit
            max_lines = max(1, int((hard_max_h - int(2 * PAD_H)) / line_h_emu))
            # Truncate text to max_lines
            all_lines = text.split("\n")
            if len(all_lines) > max_lines:
                all_lines = all_lines[:max_lines]
                all_lines[-1] = all_lines[-1][:50] + "…"
            text = "\n".join(all_lines)
            est_h = hard_max_h

        # Boundary guard
        t = self._stack_y
        if t + est_h > bottom_cap:
            est_h = bottom_cap - t
            if est_h <= 0:
                logger.warning("Stack overflow: no space for box, skipping.")
                return

        self._raw_textbox(slide, l, t, w, est_h, text, size, color, bold=bold, align=align)
        self._stack_y = t + est_h + gap_after

    # ──────────────────────────────────────────────────────────────────────────
    # RAW TEXT BOX  (lowest-level, no stacking)
    # ──────────────────────────────────────────────────────────────────────────
    def _raw_textbox(self, slide, l: int, t: int, w: int, h: int,
                     text: str, size: int, color: RGBColor,
                     *, bold: bool = False, align=PP_ALIGN.LEFT) -> None:
        """Create a text box with consistent internal padding."""
        box = slide.shapes.add_textbox(l, t, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        # Internal padding — prevents text from touching box borders
        tf.margin_top    = PAD_H
        tf.margin_bottom = PAD_H
        tf.margin_left   = PAD_W
        tf.margin_right  = PAD_W

        # Handle multi-line via separate paragraphs
        lines = text.split("\n")
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size  = Pt(size)
            p.font.bold  = bold
            p.font.name  = "Arial"
            p.font.color.rgb = color
            p.alignment  = align
            # Small paragraph spacing to prevent cramping
            p.space_after = Pt(2)

    # ──────────────────────────────────────────────────────────────────────────
    def _add_footer_box(self, slide, text: str, margin: int, width: int) -> None:
        """Storyteller insight anchored to the bottom of the slide."""
        if self._draw_frame is not None:
            fl, ft, fw, fh = self._draw_frame
            frame_bottom = ft + fh
            t = max(self._stack_y + Inches(0.08), frame_bottom - FOOTER_H - int(Inches(0.1)))
            margin = fl + int(Inches(0.05))
            width = fw - int(Inches(0.1)) - self._icon_right_pad_emu
        else:
            t = max(self._stack_y + Inches(0.1), FOOTER_T)
        bottom_cap = int(BOTTOM_SAFE)
        if self._draw_frame is not None:
            _, ft2, _, fh2 = self._draw_frame
            bottom_cap = min(bottom_cap, ft2 + fh2 - int(Inches(0.04)))
        if t + FOOTER_H > bottom_cap:
            return  # not enough room

        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, margin, t, width, FOOTER_H)
        box.fill.solid()
        box.fill.fore_color.rgb = self._light_accent(0)
        try:
            box.line.color.rgb = self._dna_color(1, 10, "highlight")
            box.line.width = Pt(0.75)
        except Exception:
            pass

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.06)
        tf.margin_bottom = Inches(0.06)
        tf.margin_left = PAD_W
        tf.margin_right = PAD_W

        p = tf.paragraphs[0]
        p.text = f"Key Insight: {clamp_text(text, max_words=20, max_chars=110)}"
        p.font.size = Pt(FS_FOOTER + 1)
        p.font.italic = True
        p.font.name = "Arial"
        p.font.color.rgb = self._dna_color(1, 10, "h1")
        dna = self.theme.design_dna
        p.alignment = PP_ALIGN.CENTER if dna.get("align") == "center" else PP_ALIGN.LEFT

    # ──────────────────────────────────────────────────────────────────────────
    def _add_icon(self, slide, hint: str, x: int, y: int) -> None:
        """Small decorative icon in brand accent color."""
        size = Inches(0.45)
        # Keep strictly within slide
        x = min(x, SW - size - Inches(0.1))
        y = max(y, Inches(0.1))
        shape_type = MSO_AUTO_SHAPE_TYPE.OVAL if hint in ("target", "shield", "globe", "users") else MSO_AUTO_SHAPE_TYPE.RECTANGLE
        icon = slide.shapes.add_shape(shape_type, x, y, size, size)
        icon.fill.solid()
        icon.fill.fore_color.rgb = self._dna_color(1, 10, "highlight")
        try:
            icon.line.color.rgb = RGBColor(255, 255, 255)
            icon.line.width = Pt(0.5)
        except Exception:
            pass

    # ──────────────────────────────────────────────────────────────────────────
    def _add_chart(self, slide, spec: dict | None, l: int, t: int, w: int, h: int) -> None:
        """Native PPTX bar chart from structured data."""
        if not spec or not spec.get("labels"):
            return
        
        labels = spec["labels"]
        values = spec["values"]
        
        data = CategoryChartData()
        data.categories = [str(lbl)[:20] for lbl in labels]
        data.add_series("Growth", [val for val in values])
        
        if t + h > BOTTOM_SAFE:
            h = BOTTOM_SAFE - t
        if h <= 0: return

        cs = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, l, t, w, h, data)
        chart = cs.chart
        chart.has_legend = False
        try:
            # High-end chart styling
            plot = chart.plots[0]
            plot.gap_width = 80 # Thick modern bars
            
            chart.font.size = Pt(10)
            chart.font.color.rgb = self._dark()
            chart.value_axis.has_major_gridlines = False
            chart.category_axis.tick_labels.font.size = Pt(9)
            
            # Use accent colors for bars
            series = chart.series[0]
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = self._accent(0)
            try:
                series.format.line.fill.background()
            except: pass
        except Exception as e:
            logger.debug("Chart styling failed: %s", e)

    @staticmethod
    def _format_metric_value(v: float | int) -> str:
        try:
            x = float(v)
        except (TypeError, ValueError):
            return str(v)
        ax = abs(x)
        if ax >= 1e12:
            return f"{x/1e12:.1f}T"
        if ax >= 1e9:
            return f"{x/1e9:.1f}B"
        if ax >= 1e6:
            return f"{x/1e6:.1f}M"
        if ax >= 1e3:
            return f"{x/1e3:.1f}K"
        if ax >= 100:
            return f"{x:,.0f}"
        return f"{x:g}"

    def _add_metric_cards(self, slide, spec: dict | None, l: int, t: int, w: int, h: int) -> None:
        """Render 3-4 numeric KPI cards with labels."""
        if not spec or not spec.get("labels"):
            return
        
        labels = spec["labels"][:4]
        values = spec["values"][:4]
        count = len(labels)
        
        card_w = (w - (count-1)*Inches(0.15)) // count
        for i in range(count):
            x = l + i * (card_w + Inches(0.15))
            
            # Rounded background
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, t, card_w, h)
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._light_accent(0)
            shape.line.color.rgb = self._primary()
            shape.line.width = Pt(1)
            
            # Value (Large)
            val_text = str(values[i])
            if isinstance(values[i], (int, float)) and values[i] >= 1000:
                val_text = self._format_metric_value(values[i])
                
            v_box = slide.shapes.add_textbox(x, t + Inches(0.2), card_w, h // 2)
            p_v = v_box.text_frame.paragraphs[0]
            p_v.text = val_text
            p_v.font.size = Pt(32)
            p_v.font.bold = True
            p_v.font.color.rgb = self._primary()
            p_v.alignment = PP_ALIGN.CENTER
            
            # Label (Small)
            l_box = slide.shapes.add_textbox(x, t + h // 2, card_w, (h // 2) - Inches(0.1))
            tf_l = l_box.text_frame
            tf_l.vertical_anchor = MSO_ANCHOR.TOP
            p_l = tf_l.paragraphs[0]
            p_l.text = str(labels[i]).upper()
            p_l.font.size = Pt(10)
            p_l.font.color.rgb = self._muted()
            p_l.alignment = PP_ALIGN.CENTER
            
            # Sub-accent bar for premium look
            bar_w = card_w // 2
            bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x + (card_w-bar_w)//2, t + h - Pt(4), bar_w, Pt(2))
            bar.fill.solid()
            bar.fill.fore_color.rgb = self._accent(i)
            bar.line.fill.background()

    def _add_icon_grid(
        self,
        slide,
        points: list[str],
        icon_paths: list[str],
        icon_tokens: list[str],
        l: int,
        t: int,
        w: int,
        h: int,
    ) -> None:
        """Render 3 vertical cards with icons and descriptions."""
        count = min(len(points), 3)
        if count == 0: return
        
        card_w = (w - (count-1)*Inches(0.2)) // count
        for i in range(count):
            x = l + i * (card_w + Inches(0.2))
            
            # Outer Card
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, t, card_w, h)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            shape.line.color.rgb = self._light_accent(0)
            shape.line.width = Pt(0.5)
            
            # Icon
            if i < len(icon_paths) and Path(icon_paths[i]).exists():
                try:
                    slide.shapes.add_picture(icon_paths[i], x + (card_w-Inches(0.6))//2, t + Inches(0.25), width=Inches(0.6))
                except Exception: pass
            elif i < len(icon_tokens):
                glyph = get_unicode_icon(icon_tokens[i])
                gbox = slide.shapes.add_textbox(x + (card_w-Inches(0.45))//2, t + Inches(0.22), Inches(0.45), Inches(0.45))
                gp = gbox.text_frame.paragraphs[0]
                gp.text = glyph
                gp.font.size = Pt(28)
                gp.font.color.rgb = self._primary()
                gp.alignment = PP_ALIGN.CENTER
            
            # Text
            txt_box = slide.shapes.add_textbox(x + Inches(0.1), t + Inches(1.0), card_w - Inches(0.2), h - Inches(1.1))
            tf = txt_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = str(points[i])
            p.font.size = Pt(11)
            p.font.color.rgb = self._dark()
            p.alignment = PP_ALIGN.CENTER

    def _add_two_col(self, slide, points: list[str], chart_spec: dict | None, l: int, t: int, w: int, h: int) -> None:
        """Left: Bullets, Right: Chart."""
        half_w = (w - Inches(0.3)) // 2
        
        # Left Bullets
        bullets = "\n".join([f"\u2022 {p}" for p in points[:4]])
        self._raw_textbox(slide, l, t, half_w, h, bullets, 11, self._dna_color(1, 10, "b"))
        
        # Right Chart
        if chart_spec:
            self._add_chart(slide, chart_spec, l + half_w + Inches(0.3), t, half_w, h)

    def _add_timeline(self, slide, points: list[str], l: int, t: int, w: int, h: int) -> None:
        """Staggered (Up/Down) premium horizontal timeline."""
        count = min(len(points), 4)
        if count == 0: return
        
        # Baseline
        line_y = t + h // 2
        line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, l, line_y - Pt(1), w, Pt(2))
        line.fill.solid()
        line.fill.fore_color.rgb = self._muted()
        line.line.fill.background()
        
        step_w = w // (count + 1)
        for i in range(count):
            x = l + (i + 1) * step_w
            is_up = i % 2 == 0
            
            # Connector vertical line
            con_h = h // 4
            con_y = line_y - con_h if is_up else line_y
            vh = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, con_y, Pt(1), con_h)
            vh.fill.solid()
            vh.fill.fore_color.rgb = self._accent(i)
            vh.line.fill.background()

            # Marker Node
            ms = Inches(0.12)
            marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x - ms//2, line_y - ms//2, ms, ms)
            marker.fill.solid()
            marker.fill.fore_color.rgb = self._accent(i)
            marker.line.color.rgb = RGBColor(255, 255, 255)
            marker.line.width = Pt(1)
            
            # Label
            box_h = h // 3
            box_y = line_y - con_h - box_h if is_up else line_y + con_h
            lbl_box = slide.shapes.add_textbox(x - step_w//2, box_y, step_w, box_h)
            tf = lbl_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = str(points[i])
            p.font.size = Pt(10)
            p.font.bold = is_up # Alternating emphasis
            p.font.color.rgb = self._dark()
            p.alignment = PP_ALIGN.CENTER

    # ──────────────────────────────────────────────────────────────────────────
    def _add_table(self, slide, tbl_spec: dict, l: int, t: int, w: int, h: int) -> None:
        """Native PPTX table with DNA styling."""
        headers = (tbl_spec.get("headers") or [])[:6]
        rows    = (tbl_spec.get("rows")    or [])[:5]
        if not headers:
            return
        num_cols = len(headers)
        num_rows = len(rows) + 1  # +1 for header

        # Boundary check
        if t + h > BOTTOM_SAFE:
            h = BOTTOM_SAFE - t
        if h <= 0:
            return

        tbl = slide.shapes.add_table(num_rows, num_cols, l, t, w, h).table

        # Header row: Bold 16pt
        for ci, hdr in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.text = str(hdr)[:25]
            self._fmt_cell(cell, fill=self._primary(), fc=RGBColor(255, 255, 255), size=FS_SUB, bold=True)

        # Data rows: 11pt
        for ri, row in enumerate(rows):
            for ci in range(num_cols):
                val = str(row[ci])[:25] if ci < len(row) else ""
                cell = tbl.cell(ri + 1, ci)
                cell.text = val
                bg = self._light_accent(0) if ri % 2 == 0 else RGBColor(250, 250, 250)
                self._fmt_cell(cell, fill=bg, fc=self._dna_color(1, 10, "b"), size=11, bold=False)

    def _fmt_cell(self, cell, *, fill: RGBColor, fc: RGBColor, size: int, bold: bool) -> None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
        tf = cell.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        # No excess padding inside cells — they're already sized
        tf.margin_top    = Inches(0.04)
        tf.margin_bottom = Inches(0.04)
        tf.margin_left   = Inches(0.08)
        tf.margin_right  = Inches(0.08)
        p = tf.paragraphs[0]
        p.font.size  = Pt(size)
        p.font.bold  = bold
        p.font.name  = "Arial"
        p.font.color.rgb = fc
        p.alignment  = PP_ALIGN.CENTER

    # ──────────────────────────────────────────────────────────────────────────
    def _add_swot(self, slide, points: list[str], l: int, t: int, w: int, h: int) -> None:
        """2×2 SWOT quadrant grid."""
        labels = ["STRENGTHS", "WEAKNESSES", "OPPORTUNITIES", "THREATS"]
        qw = (w - Inches(0.1)) // 2
        qh = (h - Inches(0.1)) // 2
        cols = [RGBColor(*self.theme.primary_color), self._accent(0),
                self._accent(1),                     self._accent(2)]
        for i in range(4):
            r, c = divmod(i, 2)
            x = l + c * (qw + Inches(0.1))
            y = t + r * (qh + Inches(0.1))
            if y + qh > BOTTOM_SAFE:
                break
            
            # Card background
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, qw, qh)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            shape.line.color.rgb = cols[i]
            shape.line.width = Pt(1.5)
            
            # Colored Header strip
            header_h = Inches(0.35)
            header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, qw, header_h)
            header.fill.solid()
            header.fill.fore_color.rgb = cols[i]
            header.line.fill.background()

            tf_hdr = header.text_frame
            p_hdr = tf_hdr.paragraphs[0]
            p_hdr.text = labels[i]
            p_hdr.font.size = Pt(12)
            p_hdr.font.bold = True
            p_hdr.font.color.rgb = self._contrast(cols[i])
            p_hdr.alignment = PP_ALIGN.CENTER

            # Content
            txt_y = y + header_h + Inches(0.05)
            txt_box = slide.shapes.add_textbox(x + Inches(0.05), txt_y, qw - Inches(0.1), qh - header_h - Inches(0.1))
            tf = txt_box.text_frame
            tf.word_wrap = True
            if i < len(points):
                p = tf.paragraphs[0]
                p.text = str(points[i])
                p.font.size = Pt(10)
                p.font.color.rgb = self._dark()
                p.alignment = PP_ALIGN.LEFT

    # ──────────────────────────────────────────────────────────────────────────
    def _add_slide_number(self, slide, current: int, total: int) -> None:
        for ph in slide.placeholders:
            if ph.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER:
                ph.text = str(current)
                return
        # Fallback: tiny text bottom-right, inside slide
        self._raw_textbox(
            slide,
            SW - Inches(1.2), SH - Inches(0.35),
            Inches(0.9), Inches(0.25),
            f"{current} / {total}", FS_FOOTER,
            self._muted(), align=PP_ALIGN.RIGHT
        )

    # ──────────────────────────────────────────────────────────────────────────
    # Internal helpers — do NOT expose old _add_textbox that bypasses stack
    # ──────────────────────────────────────────────────────────────────────────
    def _extract_chart_pairs(self, facts: list[str]) -> list[tuple[str, float]]:
        pairs = []
        for f in (facts or [])[:6]:
            m = re.search(r'([A-Za-z\s]+)[:\-]\s*([\d\.,]+)%?', str(f))
            if m:
                label = m.group(1).strip()[:20]
                try:
                    val = float(m.group(2).replace(',', ''))
                    pairs.append((label, val))
                except ValueError:
                    pass
        return pairs

    # Color helpers
    def _primary(self)       -> RGBColor: return RGBColor(*self.theme.primary_color)
    def _accent(self, i: int)-> RGBColor: return RGBColor(*self.theme.accent_colors[i % len(self.theme.accent_colors)])
    def _light_accent(self, i: int) -> RGBColor: return RGBColor(*self.theme.light_colors[i % len(self.theme.light_colors)])
    def _muted(self)         -> RGBColor: return RGBColor(*self.theme.muted_color)
    def _dark(self)          -> RGBColor: return RGBColor(*self.theme.dark_color)

    def _contrast(self, rgb: RGBColor) -> RGBColor:
        lum = (0.299 * rgb[0] + 0.587 * rgb[1] + 0.114 * rgb[2]) / 255
        return RGBColor(0, 0, 0) if lum > 0.55 else RGBColor(255, 255, 255)

    def _dna_color(self, slide_idx: int, total: int, key: str) -> RGBColor:
        dna = self.theme.design_dna
        if not dna:
            return self._dark()
        ctx = "title" if slide_idx == 0 else ("closing" if slide_idx == total - 1 else "content")
        rules = dna.get(ctx, dna.get("content", {}))
        t = rules.get(key)
        if t:
            return RGBColor(*t)
        if key == "b":        return self._dark()
        if key == "highlight": return self._accent(0)
        return self._primary()

    # ── Kept for backward compat with misc callers (table, swot) ─────────────
    def _get_dna_color(self, slide_idx: int, total_slides: int, key: str) -> RGBColor:
        return self._dna_color(slide_idx, total_slides, key)

    # ── Legacy stubs removed — these caused the duplicate _add_textbox bug ───
    def _render_default_content(self, *args, **kwargs): pass
    def _render_clean_body_points(self, *args, **kwargs): pass
    def _set_title(self, *args, **kwargs): pass
    def _add_subtitle(self, *args, **kwargs): pass
    def _get_safe_box(self, *args, **kwargs): return (Inches(0.5), Inches(1.6), Inches(12.33), Inches(4.8))
    def _add_storyteller(self, *args, **kwargs): pass

    # ── Diagrams (funnel/pyramid/chevron) ─────────────────────────────────────────────
    def _add_chevron_flow(self, slide, stages: list[str], *, left, top, width, height) -> None:
        """Draws a premium SmartArt-equivalent chevron process flow."""
        n = min(len(stages), 5) or 3
        gap = Inches(0.1)
        cw = (width - gap * (n - 1)) // n
        ch = min(height, Inches(1.5))
        cy = top + (height - ch) // 2

        for i in range(n):
            cx = left + i * (cw + gap)
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, cx, cy, cw, ch)
            shape.fill.solid()
            # Gradient-like effect by varying lightness
            shape.fill.fore_color.rgb = self._accent(i)
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.width = Pt(1.5)
            
            tf = shape.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.2) if i > 0 else Inches(0.1)
            tf.margin_right = Inches(0.2)
            p = tf.paragraphs[0]
            p.text = clamp_text(stages[i], max_words=8)
            p.font.size = Pt(FS_BODY)
            p.font.bold = True
            p.font.name = "Arial"
            p.font.color.rgb = self._contrast(self._accent(i))
            p.alignment = PP_ALIGN.CENTER

    def _add_funnel_diagram(self, slide, stages, *, left, top, width, height) -> None:
        """Draws a premium SmartArt-equivalent funnel diagram."""
        n = min(len(stages), 4) or 3
        sh = height // n
        for i in range(n):
            # Taper off dramatically like a real funnel
            wf = max(0.2, 1.0 - i * 0.25)
            sw = int(width * wf)
            x = left + (width - sw) // 2
            y = top + i * sh
            if y + sh > BOTTOM_SAFE: break
            
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.TRAPEZOID, x, y, sw, sh - Pt(2))
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._accent(i)
            # Add premium white borders
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.width = Pt(1.5)
            
            tf = shape.text_frame
            p = tf.paragraphs[0]
            # Convert text frame into vertically centered layout
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p.text = clamp_text(stages[i] if i < len(stages) else f"Stage {i+1}", max_words=8)
            p.font.size = Pt(FS_BODY + 2)
            p.font.bold = True
            p.font.name = "Arial"
            p.font.color.rgb = self._contrast(self._accent(i))
            p.alignment = PP_ALIGN.CENTER

    def _add_pyramid_diagram(self, slide, tiers, *, left, top, width, height) -> None:
        """Draws a premium SmartArt-equivalent hierarchical pyramid."""
        n = min(len(tiers), 4) or 3
        th = height // n
        for i in range(n):
            idx = n - 1 - i
            wf = (i + 1) / n
            sw = int(width * wf)
            x = left + (width - sw) // 2
            y = top + idx * th
            if y + th > BOTTOM_SAFE: break
            
            shape_type = MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE if i == 0 else MSO_AUTO_SHAPE_TYPE.TRAPEZOID
            shape = slide.shapes.add_shape(shape_type, x, y, sw, th - Pt(2))
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._light_accent(idx)
            # Add premium white borders
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.width = Pt(1.5)
            
            tf = shape.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            # Give slightly more room for text in pyramid base vs tip
            p.text = clamp_text(tiers[idx] if idx < len(tiers) else f"Tier {idx+1}", max_words=6 if i == 0 else 12)
            p.font.size = Pt(FS_BODY + (0 if i == 0 else 2))
            p.font.bold = True
            p.font.name = "Arial"
            p.font.color.rgb = self._contrast(self._light_accent(idx))
            p.alignment = PP_ALIGN.CENTER

    def _add_agenda(self, slide, points: list[str], l: int, t: int, w: int, h: int) -> None:
        """High-end numbered agenda/TOC slide."""
        count = min(len(points), 6)
        if count == 0: return
        
        row_h = h // count
        for i in range(count):
            iy = t + i * row_h
            
            # Number Circle
            ps = Inches(0.35)
            circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, l, iy + (row_h - ps)//2, ps, ps)
            circle.fill.solid()
            circle.fill.fore_color.rgb = self._primary()
            circle.line.fill.background()
            
            p_num = circle.text_frame.paragraphs[0]
            p_num.text = str(i + 1)
            p_num.font.size = Pt(14)
            p_num.font.bold = True
            p_num.font.color.rgb = RGBColor(255, 255, 255)
            p_num.alignment = PP_ALIGN.CENTER
            
            # Text box
            txt_l = l + ps + Inches(0.15)
            box = slide.shapes.add_textbox(txt_l, iy, w - (txt_l - l), row_h)
            tf = box.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = str(points[i])
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self._dark()
            p.alignment = PP_ALIGN.LEFT
            
            # Accent underline
            if i < count - 1:
                line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, txt_l, iy + row_h - Pt(1), w - (txt_l - l), Pt(0.5))
                line.fill.solid()
                line.fill.fore_color.rgb = self._light_accent(0)
                line.line.fill.background()

    def _add_key_takeaways(self, slide, points: list[str], l: int, t: int, w: int, h: int) -> None:
        """Impactful horizontal takeaway cards."""
        count = min(len(points), 3)
        if count == 0: return
        
        gap = Inches(0.2)
        cw = (w - gap * (count - 1)) // count
        
        for i in range(count):
            cx = l + i * (cw + gap)
            
            # Main Card
            card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, cx, t, cw, h)
            card.fill.solid()
            card.fill.fore_color.rgb = self._light_accent(0)
            card.line.color.rgb = self._accent(i)
            card.line.width = Pt(2)
            
            # Checkmark Icon (Simple)
            ps = Inches(0.3)
            icon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, cx + (cw - ps)//2, t - ps//2, ps, ps)
            icon.fill.solid()
            icon.fill.fore_color.rgb = self._accent(i)
            icon.line.color.rgb = RGBColor(255, 255, 255)
            
            # Text
            box = slide.shapes.add_textbox(cx + Inches(0.1), t + Inches(0.25), cw - Inches(0.2), h - Inches(0.35))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = str(points[i])
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self._primary()
            p.alignment = PP_ALIGN.CENTER

    # ── Kept to avoid import errors ───────────────────────────────────────────
    def _add_simple_chart(self, slide, title, facts, *, left, top, width, height):
        self._add_chart(slide, facts, left, top, width, height)

    def _draw_icon_strip(self, slide, tokens, icon_paths, *, left, top, width):
        if not icon_paths: return
        sz = Inches(0.6)
        sp = (width - sz * len(icon_paths)) // (len(icon_paths) + 1)
        for i, p in enumerate(icon_paths[:4]):
            if not Path(p).exists(): continue
            x = left + sp + i * (sz + sp)
            if x + sz > SW: break
            try:
                slide.shapes.add_picture(p, x, top, width=sz)
            except Exception:
                pass
