"""Fill native slide placeholders (cover / thank-you / divider) — avoids duplicating layout chrome."""
from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from md2deck.stages.layout_engine import clamp_text


def _style_paragraph(p, text: str, *, size_pt: int, color: RGBColor, bold: bool, align: PP_ALIGN) -> None:
    p.text = text
    p.font.size = Pt(size_pt)
    p.font.bold = bold
    p.font.name = "Arial"
    p.font.color.rgb = color
    p.alignment = align
    p.space_after = Pt(2)


def fill_title_subtitle_placeholders(
    slide,
    *,
    title: str,
    subtitle: str | None,
    title_rgb: RGBColor,
    subtitle_rgb: RGBColor,
    title_pt: int,
    subtitle_pt: int,
    title_align: PP_ALIGN,
    max_title_chars: int,
) -> None:
    """Populate TITLE/CENTER_TITLE and SUBTITLE placeholders; skip if none."""
    title_done = False
    sub_done = False
    t = clamp_text(title, max_words=8, max_chars=max_title_chars)
    sub = (clamp_text(subtitle, max_words=14, max_chars=120) if subtitle else "") or None

    for shape in slide.placeholders:
        try:
            pf = shape.placeholder_format
        except (AttributeError, ValueError):
            continue
        pt = pf.type
        if pt in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE) and t and not title_done:
            _style_paragraph(shape.text_frame.paragraphs[0], t, size_pt=title_pt, color=title_rgb, bold=True, align=title_align)
            title_done = True
        elif pt == PP_PLACEHOLDER.SUBTITLE and sub and not sub_done:
            _style_paragraph(shape.text_frame.paragraphs[0], sub, size_pt=subtitle_pt, color=subtitle_rgb, bold=False, align=title_align)
            sub_done = True

    # Some masters use a second TITLE-like slot; if subtitle still missing, use BODY for subtitle text
    if sub and not sub_done:
        for shape in slide.placeholders:
            try:
                pf = shape.placeholder_format
            except (AttributeError, ValueError):
                continue
            if pf.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                tf = shape.text_frame
                tf.word_wrap = True
                tf.text = sub
                p0 = tf.paragraphs[0]
                p0.font.size = Pt(subtitle_pt - 2)
                p0.font.bold = False
                p0.font.name = "Arial"
                p0.font.color.rgb = subtitle_rgb
                p0.alignment = title_align
                sub_done = True
                break


def fill_divider_placeholders(
    slide,
    *,
    title: str,
    subtitle: str | None,
    title_rgb: RGBColor,
    subtitle_rgb: RGBColor,
    title_pt: int = 24,
    subtitle_pt: int = 14,
    title_align: PP_ALIGN = PP_ALIGN.LEFT,
    max_title_chars: int = 70,
) -> None:
    fill_title_subtitle_placeholders(
        slide,
        title=title,
        subtitle=subtitle,
        title_rgb=title_rgb,
        subtitle_rgb=subtitle_rgb,
        title_pt=title_pt,
        subtitle_pt=subtitle_pt,
        title_align=title_align,
        max_title_chars=max_title_chars,
    )
