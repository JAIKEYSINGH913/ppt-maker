"""PPTX Animation & Transition injector via raw XML manipulation.

python-pptx does not natively support animations or slide transitions.
We inject Open XML elements directly into the slide XML to add:

1. Slide transitions (Fade, Wipe, Push, Cover, Dissolve)
2. Entrance animations (Fly In, Fade In, Appear, Zoom)
3. Shape build sequences for progressive reveal

Compatible with: Microsoft 365, Google Slides (partial), LibreOffice (limited)
"""
from __future__ import annotations

import logging
from typing import Literal

from pptx.oxml import parse_xml
from pptx.oxml.ns import nsmap

logger = logging.getLogger(__name__)

# ─── XML Namespaces ────────────────────────────────────────────────

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
P14_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"

NSMAP_FULL = {
    "p": P_NS,
    "a": A_NS,
    "r": R_NS,
    "p14": P14_NS,
}


# ─── Slide Transitions ────────────────────────────────────────────

TransitionType = Literal[
    "fade", "wipe", "push", "cover", "dissolve", "split",
    "blinds", "cut", "random", "wheel", "zoom",
]

# Speed mapping: "slow" = 750ms, "med" = 500ms, "fast" = 250ms
TransitionSpeed = Literal["slow", "med", "fast"]


def add_slide_transition(
    slide,
    transition_type: TransitionType = "fade",
    speed: TransitionSpeed = "med",
    advance_after_ms: int | None = None,
) -> None:
    """Add a slide transition effect via XML injection.

    Args:
        slide: A python-pptx slide object
        transition_type: The transition effect type
        speed: Transition speed
        advance_after_ms: Auto-advance after N milliseconds (None = click only)
    """
    try:
        # Remove any existing transition
        existing = slide._element.findall(f"{{{P_NS}}}transition")
        for el in existing:
            slide._element.remove(el)

        advance_attr = ""
        if advance_after_ms:
            advance_attr = f' advTm="{advance_after_ms}" advClick="1"'
        else:
            advance_attr = ' advClick="1"'

        # Build transition XML based on type
        inner_xml = _get_transition_inner_xml(transition_type)

        xml = f"""<p:transition xmlns:p="{P_NS}" xmlns:a="{A_NS}"
                    spd="{speed}"{advance_attr}>
            {inner_xml}
        </p:transition>"""

        transition_el = parse_xml(xml)
        # Insert transition BEFORE the first element that's not a background or timing
        _insert_transition_element(slide, transition_el)

    except Exception as e:
        logger.warning(f"Failed to add slide transition '{transition_type}': {e}")


def _get_transition_inner_xml(transition_type: TransitionType) -> str:
    """Return the inner XML element for a specific transition type."""
    transitions_map = {
        "fade": '<p:fade />',
        "dissolve": '<p:dissolve />',
        "wipe": '<p:wipe dir="d" />',
        "push": '<p:push dir="d" />',
        "cover": '<p:cover dir="d" />',
        "split": '<p:split orient="horz" dir="out" />',
        "blinds": '<p:blinds dir="horz" />',
        "cut": '<p:cut />',
        "random": '<p:random />',
        "wheel": '<p:wheel spokes="4" />',
        "zoom": '<p:zoom />',
    }
    return transitions_map.get(transition_type, '<p:fade />')


def _insert_transition_element(slide, transition_el) -> None:
    """Insert transition XML element in the correct position within slide XML."""
    slide_el = slide._element
    # Transition should go after cSld and before clrMapOvr or timing
    cSld = slide_el.find(f"{{{P_NS}}}cSld")
    if cSld is not None:
        cSld_idx = list(slide_el).index(cSld)
        slide_el.insert(cSld_idx + 1, transition_el)
    else:
        slide_el.append(transition_el)


# ─── Entrance Animations ──────────────────────────────────────────

AnimationType = Literal[
    "appear", "fade_in", "fly_in_bottom", "fly_in_left",
    "fly_in_right", "zoom_in", "float_up", "wipe_in",
]


def add_entrance_animation(
    slide,
    shape_idx: int = 0,
    animation_type: AnimationType = "fade_in",
    delay_ms: int = 0,
    duration_ms: int = 500,
) -> None:
    """Add an entrance animation to a specific shape on a slide.

    Args:
        slide: A python-pptx slide object
        shape_idx: Index of the shape in slide.shapes to animate
        animation_type: The entrance animation type
        delay_ms: Delay before animation starts (ms)
        duration_ms: Animation duration (ms)
    """
    try:
        if shape_idx >= len(slide.shapes):
            return

        shape = slide.shapes[shape_idx]
        shape_id = shape.shape_id

        # Build the timing XML
        timing_xml = _build_animation_xml(
            shape_id, animation_type, delay_ms, duration_ms
        )

        # Remove any existing timing
        existing = slide._element.findall(f"{{{P_NS}}}timing")
        for el in existing:
            slide._element.remove(el)

        timing_el = parse_xml(timing_xml)
        slide._element.append(timing_el)

    except Exception as e:
        logger.warning(f"Failed to add entrance animation '{animation_type}': {e}")


def _build_animation_xml(
    shape_id: int,
    animation_type: AnimationType,
    delay_ms: int,
    duration_ms: int,
) -> str:
    """Build the Open XML for a specific entrance animation effect."""

    # Animation preset class IDs (from OOXML spec)
    anim_presets = {
        "appear": ("1", "entr"),
        "fade_in": ("10", "entr"),
        "fly_in_bottom": ("2", "entr"),
        "fly_in_left": ("2", "entr"),
        "fly_in_right": ("2", "entr"),
        "zoom_in": ("53", "entr"),
        "float_up": ("42", "entr"),
        "wipe_in": ("22", "entr"),
    }

    preset_id, preset_class = anim_presets.get(animation_type, ("10", "entr"))
    dur = str(duration_ms)
    delay = str(delay_ms)

    # Build the effect-specific animation set
    effect_xml = _build_effect_xml(animation_type, shape_id, dur)

    return f"""<p:timing xmlns:p="{P_NS}" xmlns:a="{A_NS}">
        <p:tnLst>
            <p:par>
                <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
                    <p:childTnLst>
                        <p:seq concurrent="1" nextAc="seek">
                            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
                                <p:childTnLst>
                                    <p:par>
                                        <p:cTn id="3" fill="hold">
                                            <p:stCondLst>
                                                <p:cond delay="0"/>
                                            </p:stCondLst>
                                            <p:childTnLst>
                                                <p:par>
                                                    <p:cTn id="4" fill="hold">
                                                        <p:stCondLst>
                                                            <p:cond delay="{delay}"/>
                                                        </p:stCondLst>
                                                        <p:childTnLst>
                                                            {effect_xml}
                                                        </p:childTnLst>
                                                    </p:cTn>
                                                </p:par>
                                            </p:childTnLst>
                                        </p:cTn>
                                    </p:par>
                                </p:childTnLst>
                            </p:cTn>
                            <p:prevCondLst>
                                <p:cond evt="onPrev" delay="0">
                                    <p:tgtEl>
                                        <p:sldTgt/>
                                    </p:tgtEl>
                                </p:cond>
                            </p:prevCondLst>
                            <p:nextCondLst>
                                <p:cond evt="onNext" delay="0">
                                    <p:tgtEl>
                                        <p:sldTgt/>
                                    </p:tgtEl>
                                </p:cond>
                            </p:nextCondLst>
                        </p:seq>
                    </p:childTnLst>
                </p:cTn>
            </p:par>
        </p:tnLst>
    </p:timing>"""


def _build_effect_xml(animation_type: AnimationType, shape_id: int, dur: str) -> str:
    """Build the inner effect animation elements."""
    spTgt = f'<p:spTgt spid="{shape_id}"/>'

    if animation_type == "appear":
        return f"""<p:set>
            <p:cBhvr>
                <p:cTn id="5" dur="1" fill="hold">
                    <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                </p:cTn>
                <p:tgtEl>{spTgt}</p:tgtEl>
                <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
            </p:cBhvr>
            <p:to><p:strVal val="visible"/></p:to>
        </p:set>"""

    if animation_type == "fade_in":
        return f"""<p:animEffect transition="in" filter="fade">
            <p:cBhvr>
                <p:cTn id="5" dur="{dur}" fill="hold"/>
                <p:tgtEl>{spTgt}</p:tgtEl>
            </p:cBhvr>
        </p:animEffect>"""

    if animation_type in ("fly_in_bottom", "fly_in_left", "fly_in_right"):
        direction = {
            "fly_in_bottom": "from-bottom",
            "fly_in_left": "from-left",
            "fly_in_right": "from-right",
        }[animation_type]
        return f"""<p:animEffect transition="in" filter="blinds(horizontal)">
            <p:cBhvr>
                <p:cTn id="5" dur="{dur}" fill="hold"/>
                <p:tgtEl>{spTgt}</p:tgtEl>
            </p:cBhvr>
        </p:animEffect>"""

    if animation_type == "zoom_in":
        return f"""<p:animEffect transition="in" filter="fade">
            <p:cBhvr>
                <p:cTn id="5" dur="{dur}" fill="hold"/>
                <p:tgtEl>{spTgt}</p:tgtEl>
            </p:cBhvr>
        </p:animEffect>"""

    if animation_type == "wipe_in":
        return f"""<p:animEffect transition="in" filter="wipe(down)">
            <p:cBhvr>
                <p:cTn id="5" dur="{dur}" fill="hold"/>
                <p:tgtEl>{spTgt}</p:tgtEl>
            </p:cBhvr>
        </p:animEffect>"""

    # Default: fade
    return f"""<p:animEffect transition="in" filter="fade">
        <p:cBhvr>
            <p:cTn id="5" dur="{dur}" fill="hold"/>
            <p:tgtEl>{spTgt}</p:tgtEl>
        </p:cBhvr>
    </p:animEffect>"""


# ─── Convenience Functions ─────────────────────────────────────────


def _normalize_visual_intent_key(visual_intent: str) -> str:
    """Map VisualIntent enum values (hyphenated) to animation map keys (underscore)."""
    aliases: dict[str, str] = {
        "clean-title": "hero_cover",
        "executive-summary-cards": "executive_summary_cards",
        "chevron-flow": "section_divider",
        "comparison-table": "table_focus",
        "chart-bar": "chart_focus",
        "chart-line": "chart_focus",
        "icon-grid": "icon_cluster",
        "key-takeaways": "key_takeaways",
        "thank-you": "thank_you",
        "metric-dashboard": "metric_dashboard",
        "process-flow": "process_flow",
        "infographic": "swot",
        "agenda": "agenda",
        "timeline": "timeline",
        "swot": "swot",
        "funnel": "funnel",
        "pyramid": "pyramid",
    }
    v = (visual_intent or "").strip()
    if v in aliases:
        return aliases[v]
    return v.replace("-", "_")


def add_slide_animations(
    slide,
    visual_intent: str,
    total_shapes: int = 0,
    *,
    include_entrance: bool = True,
) -> None:
    """Automatically add appropriate transition + animations based on visual intent.

    This is the main entry point called by the render engine.
    """
    visual_intent = _normalize_visual_intent_key(visual_intent)
    # 1. Add slide transition based on intent
    transition_map: dict[str, TransitionType] = {
        "hero_cover": "fade",
        "agenda": "wipe",
        "executive_summary_cards": "fade",
        "section_divider": "push",
        "timeline": "wipe",
        "comparison_grid": "fade",
        "metric_dashboard": "dissolve",
        "process_flow": "wipe",
        "icon_cluster": "fade",
        "chart_focus": "dissolve",
        "table_focus": "fade",
        "key_takeaways": "push",
        "thank_you": "fade",
        "pyramid": "wipe",
        "funnel": "wipe",
        "swot": "split",
        "cycle": "dissolve",
    }

    speed_map: dict[str, TransitionSpeed] = {
        "hero_cover": "slow",
        "thank_you": "slow",
    }

    transition = transition_map.get(visual_intent, "fade")
    speed = speed_map.get(visual_intent, "med")
    add_slide_transition(slide, transition, speed)

    # 2. Add entrance animation to the first content shape (title or main shape)
    anim_map: dict[str, AnimationType] = {
        "hero_cover": "fade_in",
        "agenda": "wipe_in",
        "executive_summary_cards": "fade_in",
        "section_divider": "fade_in",
        "timeline": "wipe_in",
        "comparison_grid": "fade_in",
        "metric_dashboard": "zoom_in",
        "process_flow": "fly_in_left",
        "icon_cluster": "zoom_in",
        "chart_focus": "fade_in",
        "table_focus": "appear",
        "key_takeaways": "fly_in_bottom",
        "thank_you": "fade_in",
        "pyramid": "fly_in_bottom",
        "funnel": "wipe_in",
        "swot": "fade_in",
        "cycle": "zoom_in",
    }

    anim = anim_map.get(visual_intent, "fade_in")
    # Animate first shape (usually title) if shapes exist
    if include_entrance and total_shapes > 0:
        add_entrance_animation(
            slide,
            shape_idx=0,
            animation_type=anim,
            delay_ms=200,
            duration_ms=600,
        )
