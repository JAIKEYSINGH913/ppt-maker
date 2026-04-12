from pptx import Presentation
from pptx.util import Inches
import os

files = [
    r"c:\Users\Kaushal\Downloads\Cursor\Code EZ_ Master of Agents _ Files\Sample Files\Accenture Tech Acquisition Analysis\Template_Accenture Tech Acquisition Analysis.pptx",
    r"c:\Users\Kaushal\Downloads\Cursor\Code EZ_ Master of Agents _ Files\Sample Files\AI Bubble_ Detection, Prevention, and Investment Strategies\Template_AI Bubble_ Detection, Prevention, and Investment Strategies.pptx",
    r"c:\Users\Kaushal\Downloads\Cursor\Code EZ_ Master of Agents _ Files\Sample Files\UAE Progress toward 2050 Solar Energy Targets_20250729_120637\Template_UAE Progress toward 2050 Solar Energy Targets_20250729_120637.pptx"
]

def analyze_padding(path):
    print(f"\n--- Analyzing Padding: {os.path.basename(path)} ---")
    prs = Presentation(path)
    count = 0
    for layout in prs.slide_layouts:
        if count >= 3: break
        print(f"Layout {count}: {layout.name}")
        for ph in layout.placeholders:
            if hasattr(ph, "text_frame") and ph.text_frame:
                tf = ph.text_frame
                print(f"  Placeholder {ph.placeholder_format.type}:")
                try:
                    print(f"    Margins: L={tf.margin_left/914400.0:.2f}, T={tf.margin_top/914400.0:.2f}, R={tf.margin_right/914400.0:.2f}, B={tf.margin_bottom/914400.0:.2f}")
                except:
                    print("    Margins: Not explicitly set")
        count += 1

for f in files:
    if os.path.exists(f):
        analyze_padding(f)
