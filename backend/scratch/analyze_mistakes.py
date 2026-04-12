"""Deep analysis - Unicode safe output."""
import sys, os, itertools
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt

MISTAKES_FILE = r"c:\Users\Kaushal\Downloads\Cursor\Code EZ_ Master of Agents _ Files\Common Mistakes and overall guide to improve slides.pptx"
SAMPLE_FILES = {
    "Accenture": r"c:\Users\Kaushal\Downloads\Cursor\Code EZ_ Master of Agents _ Files\Sample Files\Accenture Tech Acquisition Analysis\Template_Accenture Tech Acquisition Analysis.pptx",
    "AI_Bubble": r"c:\Users\Kaushal\Downloads\Cursor\Code EZ_ Master of Agents _ Files\Sample Files\AI Bubble_ Detection, Prevention, and Investment Strategies\Template_AI Bubble_ Detection, Prevention, and Investment Strategies.pptx",
    "UAE_Solar": r"c:\Users\Kaushal\Downloads\Cursor\Code EZ_ Master of Agents _ Files\Sample Files\UAE Progress toward 2050 Solar Energy Targets_20250729_120637\Template_UAE Progress toward 2050 Solar Energy Targets_20250729_120637.pptx",
}

def emu_to_in(emu):
    if emu is None: return 0.0
    return round(emu / 914400, 3)

def font_size_emu_to_pt(emu_size):
    """Font size in pptx XML is in hundredths of a point (centi-points)."""
    if emu_size is None: return "inherit"
    return round(emu_size / 12700, 1)  # 12700 centi-points per pt

def analyze_pptx(path, label, max_layouts=5):
    print(f"\n{'='*60}")
    print(f"ANALYZING: {label}")
    print(f"{'='*60}")
    try:
        prs = Presentation(path)
    except Exception as e:
        print(f"  ERROR: {e}")
        return
    
    SW = emu_to_in(prs.slide_width)
    SH = emu_to_in(prs.slide_height)
    print(f"Slide: {SW}\" x {SH}\"")
    
    for mi, master in enumerate(prs.slide_masters):
        print(f"\n[Master {mi}: {len(list(master.slide_layouts))} layouts]")
        count = 0
        for layout in master.slide_layouts:
            if count >= max_layouts: break
            count += 1
            print(f"  Layout '{layout.name}':")
            phs = list(layout.placeholders)
            for ph in phs:
                l, t = emu_to_in(ph.left), emu_to_in(ph.top)
                w, h = emu_to_in(ph.width), emu_to_in(ph.height)
                ph_type = str(ph.placeholder_format.type).split('(')[0].split('.')[-1]
                flags = []
                if (l + w) > SW: flags.append("OVERFLOW_RIGHT!")
                if (t + h) > SH: flags.append("OVERFLOW_BOTTOM!")
                print(f"    PH[{ph.placeholder_format.idx}] {ph_type}: L={l}\" T={t}\" W={w}\" H={h}\" {' '.join(flags)}")
                for para in ph.text_frame.paragraphs:
                    if para.runs:
                        r = para.runs[0]
                        fsize = font_size_emu_to_pt(r.font.size)
                        print(f"      Font: {fsize}pt Bold={r.font.bold}")
                        break
    
    if prs.slides:
        print(f"\n[{len(prs.slides)} Slides]")
        for si, slide in enumerate(itertools.islice(prs.slides, 5)):
            print(f"\n  Slide {si+1}: layout='{slide.slide_layout.name}'")
            shapes_data = []
            for shape in slide.shapes:
                l, t = emu_to_in(shape.left), emu_to_in(shape.top)
                w, h = emu_to_in(shape.width), emu_to_in(shape.height)
                flags = []
                if (l + w) > SW: flags.append("RIGHT_OVERFLOW!")
                if (t + h) > SH: flags.append("BOTTOM_OVERFLOW!")
                txt = ""
                if hasattr(shape, "text"):
                    txt = repr(shape.text[:50])
                fsize_info = ""
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        if para.runs:
                            r = para.runs[0]
                            fsize_info = f"{font_size_emu_to_pt(r.font.size)}pt"
                            break
                print(f"    [{shape.shape_type}] L={l}\" T={t}\" W={w}\" H={h}\" {fsize_info} {txt} {' '.join(flags)}")
                shapes_data.append((l, t, w, h))
            for i in range(len(shapes_data)):
                for j in range(i+1, len(shapes_data)):
                    al, at, aw, ah = shapes_data[i]
                    bl, bt, bw, bh = shapes_data[j]
                    if al < (bl+bw) and (al+aw) > bl and at < (bt+bh) and (at+ah) > bt:
                        print(f"    !! SHAPE OVERLAP: shape[{i}] vs shape[{j}]")

analyze_pptx(MISTAKES_FILE, "Common Mistakes Guide", max_layouts=1)
for label, path in SAMPLE_FILES.items():
    if os.path.exists(path):
        analyze_pptx(path, label, max_layouts=4)
    else:
        print(f"\nNot found: {path}")
