from pptx import Presentation
from pptx.util import Inches
import os

files = [
    r"c:\Users\Kaushal\Downloads\Cursor\Code EZ_ Master of Agents _ Files\Sample Files\Accenture Tech Acquisition Analysis\Template_Accenture Tech Acquisition Analysis.pptx",
    r"c:\Users\Kaushal\Downloads\Cursor\Code EZ_ Master of Agents _ Files\Sample Files\AI Bubble_ Detection, Prevention, and Investment Strategies\Template_AI Bubble_ Detection, Prevention, and Investment Strategies.pptx",
    r"c:\Users\Kaushal\Downloads\Cursor\Code EZ_ Master of Agents _ Files\Sample Files\UAE Progress toward 2050 Solar Energy Targets_20250729_120637\Template_UAE Progress toward 2050 Solar Energy Targets_20250729_120637.pptx"
]

def analyze_pptx(path):
    print(f"\n--- Analyzing: {os.path.basename(path)} ---")
    prs = Presentation(path)
    print(f"Slide Size: {prs.slide_width / 914400.0:.2f} x {prs.slide_height / 914400.0:.2f} inches")
    
    # Analyze common layouts
    for l_idx in [0, 1, 4]:
        if l_idx >= len(prs.slide_layouts): continue
        layout = prs.slide_layouts[l_idx]
        print(f"\n[Layout {l_idx}: {layout.name}]")
        
        for ph in layout.placeholders:
            pt = ph.placeholder_format.type
            print(f"  Placeholder: {pt} (idx={ph.placeholder_format.idx})")
            print(f"    Pos: L={ph.left/914400.0:.2f}, T={ph.top/914400.0:.2f}")
            
            if hasattr(ph, "text_frame") and ph.text_frame.paragraphs:
                p = ph.text_frame.paragraphs[0]
                if p.font.size: print(f"    Font Size: {p.font.size.pt:.1f}pt")
                if p.font.color and hasattr(p.font.color, 'rgb'): print(f"    Text Color: {p.font.color.rgb}")

        # Deep dive into shapes (Icons/Infographics) and Charts
        for shape in layout.shapes:
            if shape.is_placeholder: continue
            print(f"  Shape: {shape.name} ({shape.shape_type})")
            
            # Text within shapes?
            if hasattr(shape, "text_frame") and shape.text_frame.paragraphs:
                p = shape.text_frame.paragraphs[0]
                if p.font.color and hasattr(p.font.color, 'rgb'):
                    print(f"    Shape Text Color: {p.font.color.rgb}")

            # Fill colors
            if hasattr(shape, "fill"):
                from pptx.enum.dml import MSO_FILL
                if shape.fill.type == MSO_FILL.SOLID:
                    print(f"    Solid Fill: {shape.fill.fore_color.rgb if hasattr(shape.fill.fore_color, 'rgb') else 'N/A'}")
                elif shape.fill.type == MSO_FILL.GRADIENT:
                    print(f"    Gradient detected")
            
            if shape.has_chart:
                chart = shape.chart
                print(f"    Chart Type: {chart.chart_type}")

for f in files:
    if os.path.exists(f):
        analyze_pptx(f)
    else:
        print(f"Not found: {f}")
